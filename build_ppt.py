from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
MAROON   = RGBColor(0x80, 0x00, 0x00)
NAVY     = RGBColor(0x1a, 0x1a, 0x2e)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY    = RGBColor(0xF2, 0xF2, 0xF2)
DGRAY    = RGBColor(0x33, 0x33, 0x33)
MGRAY    = RGBColor(0x88, 0x88, 0x88)
CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG  = RGBColor(0xA8, 0xE6, 0xA3)
ORANGE   = RGBColor(0xE8, 0x7A, 0x1E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h,
                font_size=Pt(18), bold=False, color=DGRAY,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def header_bar(slide, title, subtitle=None):
    """Maroon top bar with white title text."""
    bar = add_rect(slide, 0, 0, W, Inches(1.25), fill=MAROON)
    # title text
    txb = slide.shapes.add_textbox(Inches(0.4), Inches(0.12), W - Inches(0.8), Inches(0.9))
    txb.word_wrap = False
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(15)
        r2.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
    return bar

def slide_bg(slide, color=WHITE):
    add_rect(slide, 0, 0, W, H, fill=color)

def bullet_block(slide, items, l, t, w, h,
                 font_size=Pt(20), color=DGRAY, line_spacing=1.3, icon="▸"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = font_size
        run.font.color.rgb = color

def code_block(slide, code, l, t, w, h, font_size=Pt(15)):
    bg = add_rect(slide, l, t, w, h, fill=CODE_BG)
    txb = slide.shapes.add_textbox(l + Inches(0.2), t + Inches(0.15),
                                   w - Inches(0.4), h - Inches(0.3))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in code.strip().split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = font_size
        run.font.name = "Courier New"
        run.font.color.rgb = CODE_FG

def pill(slide, text, l, t, w=Inches(2.2), h=Inches(0.55),
         fill=MAROON, font_color=WHITE, font_size=Pt(16)):
    add_rect(slide, l, t, w, h, fill=fill)
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = True
    run.font.color.rgb = font_color

def add_table(slide, headers, rows, l, t, w, h,
              hdr_fill=MAROON, hdr_fg=WHITE,
              row_fills=None, font_size=Pt(16)):
    cols = len(headers)
    nrows = len(rows) + 1
    tbl = slide.shapes.add_table(nrows, cols, l, t, w, h).table
    col_w = w // cols
    for i in range(cols):
        tbl.columns[i].width = col_w

    def set_cell(cell, text, bold=False, fill=None, fg=DGRAY, align=PP_ALIGN.LEFT):
        cell.text = text
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = fg
        if fill:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill

    for ci, hdr in enumerate(headers):
        set_cell(tbl.cell(0, ci), hdr, bold=True, fill=hdr_fill, fg=hdr_fg,
                 align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        rf = (row_fills[ri] if row_fills and ri < len(row_fills)
              else (LGRAY if ri % 2 == 0 else WHITE))
        for ci, val in enumerate(row):
            set_cell(tbl.cell(ri+1, ci), val, fill=rf, fg=DGRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
# Navy background
add_rect(sl, 0, 0, W, H, fill=NAVY)
# Maroon accent strip on left
add_rect(sl, 0, 0, Inches(0.35), H, fill=MAROON)

add_textbox(sl, "Adding a Shiny App to the\nCRI Biocore Platform",
            Inches(0.8), Inches(1.6), Inches(10), Inches(2.4),
            font_size=Pt(48), bold=True, color=WHITE)

add_textbox(sl, "A step-by-step guide  ·  Git  →  Claude Code  →  GitHub Actions  →  ShinyProxy",
            Inches(0.8), Inches(4.2), Inches(11), Inches(0.7),
            font_size=Pt(22), color=RGBColor(0xCC, 0xCC, 0xCC))

add_textbox(sl, "CRI Bioinformatics Core",
            Inches(0.8), Inches(6.5), Inches(5), Inches(0.6),
            font_size=Pt(16), color=RGBColor(0x99, 0x99, 0x99))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Big Picture (pipeline diagram)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "The Big Picture", "End-to-end workflow overview")

STEPS = [
    ("1", "Your R\nApp Files", "(local)", MAROON),
    ("2", "GitHub\nRepo",    "(PR)",    RGBColor(0x1565,0xC0,0x00) if False else RGBColor(0x15,0x65,0xC0)),
    ("3", "GitHub\nActions", "(auto)",  RGBColor(0xE6,0x7E,0x22)),
    ("4", "ShinyProxy\nServer", "(deploy)", RGBColor(0x27,0xAE,0x60)),
    ("5", "Live App\nin Browser", "(done!)",  RGBColor(0x6A,0x1B,0x9A)),
]

box_w = Inches(2.0)
box_h = Inches(1.8)
gap   = Inches(0.15)
arrow_w = Inches(0.4)
total = len(STEPS) * box_w + (len(STEPS)-1) * (gap + arrow_w)
start_x = (W - total) / 2
y = Inches(2.8)

for i, (num, label, sub, color) in enumerate(STEPS):
    x = start_x + i * (box_w + gap + arrow_w)
    # box
    add_rect(sl, x, y, box_w, box_h, fill=color)
    # number circle area
    add_textbox(sl, num, x, y + Inches(0.1), box_w, Inches(0.5),
                font_size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, label, x, y + Inches(0.55), box_w, Inches(0.85),
                font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, sub, x, y + Inches(1.35), box_w, Inches(0.4),
                font_size=Pt(13), color=RGBColor(0xFF,0xEE,0xEE), align=PP_ALIGN.CENTER)
    # arrow (except after last)
    if i < len(STEPS)-1:
        ax = x + box_w + gap
        add_textbox(sl, "→", ax, y + Inches(0.55), arrow_w, box_h - Inches(0.55),
                    font_size=Pt(28), bold=True, color=MGRAY, align=PP_ALIGN.CENTER)

add_textbox(sl, "You write R code.  Claude + GitHub Actions handle everything else.",
            Inches(1), Inches(5.4), Inches(11), Inches(0.55),
            font_size=Pt(20), color=MGRAY, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Prerequisites
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "What You Need Before You Start")

items = [
    "Your Shiny app files  (app.R  or  server.R + ui.R)",
    "A GitHub account with access to the CRI Biocore Apps repo",
    "Claude Code CLI installed  →  npm install -g @anthropic-ai/claude-code",
    "No Docker knowledge required — Claude generates the Dockerfile for you",
    "No manual YAML editing — Claude adds every config entry automatically",
]
bullet_block(sl, items,
             Inches(0.6), Inches(1.5), Inches(12), Inches(5),
             font_size=Pt(22))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Step 1: Clone & drop files
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Step 1 — Clone the Repo & Add Your App")

# Left column — text
add_textbox(sl, "1. Clone the repository",
            Inches(0.5), Inches(1.45), Inches(5.8), Inches(0.5),
            font_size=Pt(19), bold=True, color=MAROON)
code_block(sl,
           "git clone https://github.com/\n"
           "  zhongyuli2026/CRI_biocore_apps.git\n\n"
           "cd CRI_biocore_apps",
           Inches(0.5), Inches(1.95), Inches(5.8), Inches(1.6))

add_textbox(sl, "2. Drop your app folder inside  apps/",
            Inches(0.5), Inches(3.65), Inches(5.8), Inches(0.5),
            font_size=Pt(19), bold=True, color=MAROON)

bullet_block(sl,
             ["Folder name uses hyphens  (e.g.  030-basic-datatable)",
              "Put all .R files in that folder",
              "No Dockerfile needed yet — Claude writes it"],
             Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.2),
             font_size=Pt(18))

# Right column — folder tree visual
add_rect(sl, Inches(7.0), Inches(1.4), Inches(5.8), Inches(5.5), fill=CODE_BG)
tree = (
    "CRI_biocore_apps/\n"
    "├── apps/\n"
    "│   ├── DEApp-master/\n"
    "│   ├── scvizapp/\n"
    "│   └── 030-basic-datatable/   ← your app\n"
    "│       ├── server.R\n"
    "│       └── ui.R\n"
    "├── shinyproxy/\n"
    "│   ├── application.yml\n"
    "│   └── application-public.yml\n"
    "└── .github/\n"
    "    └── workflows/\n"
    "        └── deploy.yml"
)
txb = sl.shapes.add_textbox(Inches(7.2), Inches(1.6), Inches(5.4), Inches(5.1))
txb.word_wrap = False
tf = txb.text_frame
tf.word_wrap = False
first = True
for line in tree.split("\n"):
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.size = Pt(14)
    run.font.name = "Courier New"
    run.font.color.rgb = CODE_FG


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Step 2: Claude /add-shiny-app
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Step 2 — Run the Claude Skill")

add_textbox(sl, "Open Claude Code in the repo folder and run:",
            Inches(0.6), Inches(1.45), Inches(12), Inches(0.5),
            font_size=Pt(20), color=DGRAY)
code_block(sl, "claude\n/add-shiny-app",
           Inches(0.6), Inches(2.0), Inches(5.5), Inches(1.3), font_size=Pt(22))

add_textbox(sl, "Claude will interactively ask you:",
            Inches(0.6), Inches(3.5), Inches(12), Inches(0.45),
            font_size=Pt(20), bold=True, color=MAROON)

questions = [
    "App name?  (becomes the folder name, image name, and app ID)",
    "Public or private?  (controls which config file it's registered in)",
    "Display name for the dashboard?",
    "Confirm the R packages it detected?",
]
bullet_block(sl, questions,
             Inches(0.6), Inches(4.0), Inches(12.2), Inches(3.0),
             font_size=Pt(20), icon="❓")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — What Claude Generates
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "What Claude Generates — Automatically")

files = [
    ("Dockerfile",
     "apps/your-app/",
     "Packages your R app + all dependencies into a Docker container"),
    ("application-public.yml\n(or application.yml)",
     "shinyproxy/",
     "Registers your app on the server with display name, description, and image"),
    ("deploy.yml",
     ".github/workflows/",
     "Tells GitHub Actions to build the image and pull it to the server on every push"),
]

for i, (fname, path, desc) in enumerate(files):
    y = Inches(1.55) + i * Inches(1.85)
    add_rect(sl, Inches(0.4), y, Inches(12.5), Inches(1.65),
             fill=LGRAY if i % 2 == 0 else WHITE)
    # icon box
    add_rect(sl, Inches(0.4), y, Inches(0.15), Inches(1.65), fill=MAROON)
    add_textbox(sl, fname, Inches(0.65), y + Inches(0.1), Inches(3.5), Inches(0.8),
                font_size=Pt(18), bold=True, color=MAROON)
    add_textbox(sl, path, Inches(0.65), y + Inches(0.85), Inches(3.5), Inches(0.6),
                font_size=Pt(14), color=MGRAY, italic=True)
    add_textbox(sl, desc, Inches(4.4), y + Inches(0.4), Inches(8.5), Inches(0.9),
                font_size=Pt(17), color=DGRAY)

add_textbox(sl,
            "Claude scans your library() calls, resolves system dependencies, and fills in all names consistently.",
            Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
            font_size=Pt(15), color=MGRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Closer look: Dockerfile
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "A Closer Look — The Dockerfile", "Think of it as a recipe for your app's environment")

add_textbox(sl, "Claude auto-detects your R packages and writes this for you:",
            Inches(0.6), Inches(1.4), Inches(12), Inches(0.45),
            font_size=Pt(18), color=DGRAY)

dockerfile = """\
FROM rocker/shiny:4.4.2          # Start with R + Shiny pre-installed

# Install R packages detected from your library() calls
RUN R -e "install.packages(c('ggplot2', 'DT'))"

# Copy your app files into the container
RUN rm -rf /srv/shiny-server/*
COPY . /srv/shiny-server/

EXPOSE 3838
CMD ["/usr/bin/shiny-server"]"""

code_block(sl, dockerfile,
           Inches(0.6), Inches(1.95), Inches(7.8), Inches(4.3))

# Annotation boxes on right
notes = [
    (Inches(2.05), "Base image: R 4.4.2\n+ Shiny Server"),
    (Inches(2.95), "Packages auto-\ndetected by Claude"),
    (Inches(4.6),  "App files copied in"),
    (Inches(5.45), "Standard Shiny port"),
]
for (ny, note) in notes:
    add_rect(sl, Inches(8.7), ny, Inches(4.2), Inches(0.75),
             fill=RGBColor(0xFF,0xF3,0xE0), line=ORANGE, line_width=Pt(1))
    add_textbox(sl, note, Inches(8.85), ny + Inches(0.05), Inches(3.9), Inches(0.65),
                font_size=Pt(14), color=RGBColor(0x8B,0x45,0x13))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Step 3: Commit & PR
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Step 3 — Commit and Open a Pull Request")

add_textbox(sl, "Create a branch, commit all changes, and push:",
            Inches(0.6), Inches(1.4), Inches(12), Inches(0.45),
            font_size=Pt(19), color=DGRAY)

code_block(sl,
           "git checkout -b yourname\n\n"
           "git add apps/your-app/Dockerfile \\\n"
           "        shinyproxy/application-public.yml \\\n"
           "        .github/workflows/deploy.yml\n\n"
           "git commit -m \"add 030-basic-datatable\"\n\n"
           "git push origin yourname",
           Inches(0.6), Inches(1.95), Inches(6.2), Inches(4.5))

add_textbox(sl, "Then on GitHub:",
            Inches(7.2), Inches(1.95), Inches(5.7), Inches(0.5),
            font_size=Pt(19), bold=True, color=MAROON)

steps = [
    "Open a Pull Request targeting main",
    "Get it reviewed and approved",
    "Merge the PR  →  deployment starts automatically",
]
bullet_block(sl, steps,
             Inches(7.2), Inches(2.5), Inches(5.7), Inches(2.0),
             font_size=Pt(18))

# Warning box
add_rect(sl, Inches(7.2), Inches(5.0), Inches(5.7), Inches(1.3),
         fill=RGBColor(0xFF,0xF0,0xF0), line=MAROON, line_width=Pt(2))
add_textbox(sl, "⚠   Never push directly to main",
            Inches(7.35), Inches(5.1), Inches(5.4), Inches(0.5),
            font_size=Pt(17), bold=True, color=MAROON)
add_textbox(sl, "Pushing to main triggers an immediate\nproduction deployment.",
            Inches(7.35), Inches(5.55), Inches(5.4), Inches(0.65),
            font_size=Pt(15), color=DGRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — What is GitHub Actions?
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "What is GitHub Actions?", "A built-in automation engine inside GitHub")

# Left: plain-language description
add_textbox(sl, "The short version:",
            Inches(0.6), Inches(1.45), Inches(5.8), Inches(0.45),
            font_size=Pt(19), bold=True, color=MAROON)
add_textbox(sl,
            "GitHub Actions is a robot that watches your\nrepository and automatically runs tasks\nwhenever something happens — like a PR merge.",
            Inches(0.6), Inches(1.95), Inches(5.8), Inches(1.5),
            font_size=Pt(19), color=DGRAY)

add_textbox(sl, "Key concepts:",
            Inches(0.6), Inches(3.55), Inches(5.8), Inches(0.45),
            font_size=Pt(19), bold=True, color=MAROON)
concepts = [
    "Workflow — a script written in YAML that defines the tasks",
    "Trigger — the event that starts the workflow (e.g. push to main)",
    "Runner — a machine (cloud or self-hosted) that executes the steps",
    "Step — one command or action inside the workflow",
]
bullet_block(sl, concepts,
             Inches(0.6), Inches(4.05), Inches(5.8), Inches(2.8),
             font_size=Pt(17))

# Right: mini workflow YAML example
add_textbox(sl, "Our deploy.yml in plain English:",
            Inches(7.0), Inches(1.45), Inches(5.9), Inches(0.45),
            font_size=Pt(17), bold=True, color=MAROON)
code_block(sl,
           "on: push to main branch\n\n"
           "jobs:\n"
           "  build:\n"
           "    - Check out the code\n"
           "    - Build Docker image\n"
           "    - Push image to GHCR\n\n"
           "  deploy:\n"
           "    - SSH into production server\n"
           "    - Pull new Docker image\n"
           "    - Restart ShinyProxy",
           Inches(7.0), Inches(1.95), Inches(5.9), Inches(4.0),
           font_size=Pt(16))

add_textbox(sl,
            "The workflow file lives at  .github/workflows/deploy.yml  and is managed by Claude.",
            Inches(0.6), Inches(7.0), Inches(12.3), Inches(0.38),
            font_size=Pt(15), color=MGRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — GitHub Actions (workflow detail — existing)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Step 4 — GitHub Actions Takes Over (Automatically)")

add_textbox(sl, "Triggered by: merging the PR into main",
            Inches(0.6), Inches(1.4), Inches(12), Inches(0.45),
            font_size=Pt(19), color=MGRAY, italic=True)

actions = [
    ("Check out", "Downloads the latest code from the repository"),
    ("Build image", "Builds a Docker image from your Dockerfile on GitHub's servers"),
    ("Push to GHCR", "Pushes the image to GitHub Container Registry (ghcr.io)"),
    ("SSH to server", "Connects to the production server via a self-hosted runner"),
    ("Pull & restart", "Pulls the new image and restarts ShinyProxy — app is live!"),
]

for i, (step, desc) in enumerate(actions):
    y = Inches(2.0) + i * Inches(0.98)
    # number pill
    add_rect(sl, Inches(0.5), y + Inches(0.1), Inches(0.55), Inches(0.65), fill=MAROON)
    add_textbox(sl, str(i+1), Inches(0.5), y + Inches(0.1), Inches(0.55), Inches(0.65),
                font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(sl, step, Inches(1.2), y + Inches(0.1), Inches(2.5), Inches(0.65),
                font_size=Pt(18), bold=True, color=MAROON)
    add_textbox(sl, desc, Inches(3.9), y + Inches(0.1), Inches(9.0), Inches(0.65),
                font_size=Pt(18), color=DGRAY)
    if i < len(actions)-1:
        add_textbox(sl, "↓", Inches(0.6), y + Inches(0.72), Inches(0.4), Inches(0.28),
                    font_size=Pt(14), color=MGRAY, align=PP_ALIGN.CENTER)

add_textbox(sl, "Your only job: merge the PR.",
            Inches(0.5), Inches(7.0), Inches(12), Inches(0.38),
            font_size=Pt(17), bold=True, color=MAROON, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — What is ShinyProxy?
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "What is ShinyProxy?", "The platform that serves Shiny apps to multiple users")

# Left column — explanation
add_textbox(sl, "The problem it solves:",
            Inches(0.6), Inches(1.45), Inches(5.8), Inches(0.45),
            font_size=Pt(19), bold=True, color=MAROON)
add_textbox(sl,
            "A basic Shiny Server shares one R process\namong all users — they interfere with each other.\nShinyProxy gives every user their own\nisolated container.",
            Inches(0.6), Inches(1.95), Inches(5.8), Inches(1.8),
            font_size=Pt(18), color=DGRAY)

add_textbox(sl, "What ShinyProxy does:",
            Inches(0.6), Inches(3.85), Inches(5.8), Inches(0.45),
            font_size=Pt(19), bold=True, color=MAROON)
features = [
    "Reads application.yml to know which apps to serve",
    "Handles user login (CNet / LDAP authentication)",
    "Spins up a fresh Docker container per user session",
    "Stops the container when the user closes the app",
    "Supports public (no login) and private apps side by side",
]
bullet_block(sl, features,
             Inches(0.6), Inches(4.35), Inches(5.8), Inches(2.8),
             font_size=Pt(17))

# Right column — comparison table
add_textbox(sl, "Basic Shiny Server  vs  ShinyProxy:",
            Inches(7.0), Inches(1.45), Inches(5.9), Inches(0.45),
            font_size=Pt(17), bold=True, color=MAROON)
cmp_headers = ["Feature", "Shiny Server", "ShinyProxy"]
cmp_rows = [
    ["User isolation",    "❌ Shared",      "✅ Per container"],
    ["Authentication",    "❌ None",         "✅ LDAP / CNet"],
    ["Scaling",           "Limited",         "✅ Docker-based"],
    ["Config",            "config file",     "application.yml"],
    ["Used by CRI",       "—",               "✅ Yes"],
]
add_table(sl, cmp_headers, cmp_rows,
          Inches(7.0), Inches(2.0), Inches(5.9), Inches(4.5),
          font_size=Pt(15))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ShinyProxy diagram
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "How ShinyProxy Serves Your App")

# Simple diagram: browser → ShinyProxy → containers
elements = [
    (Inches(0.5),  Inches(3.0), RGBColor(0x15,0x65,0xC0), "User\nBrowser"),
    (Inches(4.8),  Inches(3.0), MAROON,                    "ShinyProxy\nServer"),
    (Inches(9.0),  Inches(2.0), RGBColor(0x27,0xAE,0x60),  "App A\nContainer"),
    (Inches(9.0),  Inches(3.6), RGBColor(0x27,0xAE,0x60),  "App B\nContainer"),
    (Inches(9.0),  Inches(5.2), RGBColor(0x27,0xAE,0x60),  "App C\nContainer"),
]
bw, bh = Inches(2.8), Inches(1.3)
for (bx, by, bc, bt) in elements:
    add_rect(sl, bx, by, bw, bh, fill=bc)
    add_textbox(sl, bt, bx, by + Inches(0.2), bw, bh - Inches(0.2),
                font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrows
add_textbox(sl, "→", Inches(3.5), Inches(3.35), Inches(1.1), Inches(0.6),
            font_size=Pt(28), bold=True, color=MGRAY, align=PP_ALIGN.CENTER)
add_textbox(sl, "→", Inches(7.8), Inches(2.35), Inches(1.1), Inches(0.6),
            font_size=Pt(28), bold=True, color=MGRAY, align=PP_ALIGN.CENTER)
add_textbox(sl, "→", Inches(7.8), Inches(3.95), Inches(1.1), Inches(0.6),
            font_size=Pt(28), bold=True, color=MGRAY, align=PP_ALIGN.CENTER)
add_textbox(sl, "→", Inches(7.8), Inches(5.55), Inches(1.1), Inches(0.6),
            font_size=Pt(28), bold=True, color=MGRAY, align=PP_ALIGN.CENTER)

add_textbox(sl, "reads application.yml\nto know which apps exist",
            Inches(4.8), Inches(4.45), Inches(2.8), Inches(0.7),
            font_size=Pt(12), color=MGRAY, italic=True, align=PP_ALIGN.CENTER)

bullets = [
    "Each user gets their own isolated Docker container session",
    "Container starts when user opens the app, stops when they close it",
    "Public apps: no login  ·  Private apps: CNet (LDAP) authentication",
]
bullet_block(sl, bullets,
             Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.3),
             font_size=Pt(17), color=DGRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — End Result
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl, color=NAVY)
add_rect(sl, 0, 0, Inches(0.35), H, fill=MAROON)

add_textbox(sl, "Your App Is Live",
            Inches(0.8), Inches(1.0), Inches(11), Inches(1.0),
            font_size=Pt(48), bold=True, color=WHITE)

results = [
    "Accessible at the CRI Biocore Apps dashboard",
    "Public apps: open to anyone · Private apps: CNet login required",
    "Thumbnail and description shown on the landing page",
    "Automatically rebuilt and redeployed on every future PR merge",
]
bullet_block(sl, results,
             Inches(0.8), Inches(2.3), Inches(11.5), Inches(3.5),
             font_size=Pt(24), color=WHITE)

add_textbox(sl,
            "Need to update your app?  Just edit the code, open a new PR, and merge.\nEverything else is automatic.",
            Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.0),
            font_size=Pt(18), color=RGBColor(0xCC,0xCC,0xCC), italic=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Summary
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Summary — The 4-Step Workflow")

headers = ["Step", "Action", "Who Does It", "Time"]
rows = [
    ["1", "Clone repo and drop your app files into apps/", "You", "~1 min"],
    ["2", "Run  /add-shiny-app  in Claude Code", "Claude", "~3 min"],
    ["3", "Commit, push, open a Pull Request", "You", "~1 min"],
    ["4", "Merge PR → build → deploy", "GitHub Actions", "~5 min auto"],
]
add_table(sl, headers, rows,
          Inches(0.5), Inches(1.5), Inches(12.3), Inches(3.6),
          font_size=Pt(18))

add_textbox(sl, "Total hands-on effort:  ~5 minutes",
            Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.55),
            font_size=Pt(22), bold=True, color=MAROON, align=PP_ALIGN.CENTER)

add_textbox(sl,
            "Questions?  Contact the CRI Bioinformatics Core  ·  cri.uchicago.edu/bioinformatics",
            Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
            font_size=Pt(16), color=MGRAY, align=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────────────
out = "/Users/zhongyu/CRI_biocore_apps/CRI_Biocore_Shiny_Workflow.pptx"
prs.save(out)
print(f"Saved: {out}")
